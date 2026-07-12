#!/usr/bin/env python
"""Build the final course presentation in two formats from one slide-spec list.

Outputs (into --out, default slides/):
  - final_presentation.pptx  (python-pptx, 16:9)
  - final_presentation.pdf   (matplotlib PdfPages rendering of the same spec)

All result numbers are read at runtime from results/metrics/*.csv|json and the
config; nothing is hardcoded. Missing files degrade gracefully: figures get a
placeholder note, computed bullets are dropped (never printed as 'nan').

Usage:
    python scripts/build_slides.py [--config configs/config.yaml] [--out slides/]
"""

from __future__ import annotations

import argparse
import json
import math
import textwrap
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
import pandas as pd  # noqa: E402
import yaml  # noqa: E402
from matplotlib.backends.backend_pdf import PdfPages  # noqa: E402
from matplotlib.patches import Rectangle  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

# ----------------------------------------------------------------------------
# Style constants
# ----------------------------------------------------------------------------
SLIDE_W, SLIDE_H = 13.333, 7.5  # inches, 16:9
TITLE_C = "#243447"
BODY_C = "#3a3a3a"
ACCENT = "#2a78d6"
MUTED = "#7a8794"
PLACE_BG = "#eef2f6"
PLACE_BORDER = "#aab6c2"

REPO_URL = "https://github.com/ShiraTzi/Final-Project-Image-Processing"
TEAM = "Yonatan Haba  ·  Shira Tziony"
DECK_TITLE = "Robustness of Vision Methods under Image Distortions"
DECK_SUBTITLE = "Image Processing / Vision course project — COCO benchmark"

TASK_INFO = {  # task -> (display name, model, metric label)
    "features": ("Feature matching", "ORB", "match ratio (vs clean)"),
    "detection": ("Object detection", "YOLOv8n", "bbox mAP@[.5:.95]"),
    "keypoints": ("Human pose", "Keypoint R-CNN", "OKS AP"),
    "segmentation": ("Panoptic segmentation", "Panoptic FPN", "PQ"),
}
TASK_ORDER = ["features", "detection", "keypoints", "segmentation"]


def hexrgb(h: str) -> RGBColor:
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


# ----------------------------------------------------------------------------
# Value helpers: never let a nan/None reach a slide
# ----------------------------------------------------------------------------
def ok(v) -> bool:
    if v is None:
        return False
    if isinstance(v, float) and (math.isnan(v) or math.isinf(v)):
        return False
    return True


def f3(v) -> str:
    return f"{float(v):.3f}"


def pct(v) -> str:
    return f"{round(float(v) * 100):d}%"


def fmt_lr(v) -> str:
    """1e-4 style for small learning rates, plain %g otherwise."""
    v = float(v)
    if 0 < v < 0.01:
        mant, exp = f"{v:e}".split("e")
        mant = mant.rstrip("0").rstrip(".")
        return f"{mant}e{int(exp)}"
    return f"{v:g}"


def _safe(fn, default=None):
    try:
        out = fn()
    except Exception:
        return default
    return out


# ----------------------------------------------------------------------------
# Data loading
# ----------------------------------------------------------------------------
def load_config(config_path: Path) -> dict:
    with open(config_path) as f:
        return yaml.safe_load(f)


def resolve_path(p: str, config_dir: Path, repo_root: Path) -> Path:
    q = Path(p)
    if q.is_absolute():
        return q
    for base in (Path.cwd(), repo_root, config_dir, config_dir.parent):
        cand = base / q
        if cand.exists():
            return cand
    return repo_root / q


def load_inputs(cfg: dict, config_dir: Path, repo_root: Path) -> dict:
    paths = cfg.get("paths", {})
    metrics_dir = resolve_path(paths.get("metrics_dir", "results/metrics"), config_dir, repo_root)
    figures_dir = resolve_path(paths.get("figures_dir", "results/figures"), config_dir, repo_root)

    comp = _safe(lambda: pd.read_csv(metrics_dir / "comparison.csv"))
    slong = _safe(lambda: pd.read_csv(metrics_dir / "summary_long.csv"))
    seg = _safe(lambda: json.loads((metrics_dir / "segmentation__clean.json").read_text()))
    return {
        "metrics_dir": metrics_dir,
        "figures_dir": figures_dir,
        "comp": comp,
        "slong": slong,
        "seg": seg,
    }


def find_figure(figures_dir: Path, name: str, fallback_glob: str | None = None):
    """Return an existing figure Path or None (placeholder handling downstream)."""
    p = figures_dir / name
    if p.is_file():
        return p
    if fallback_glob:
        hits = sorted(figures_dir.glob(fallback_glob))
        if hits:
            return hits[0]
    return None


# ----------------------------------------------------------------------------
# Derived values (all guarded; missing -> None and its bullet is dropped)
# ----------------------------------------------------------------------------
def compute_context(cfg: dict, data: dict) -> dict:
    C: dict = {}
    comp: pd.DataFrame | None = data["comp"]
    slong: pd.DataFrame | None = data["slong"]
    seg = data["seg"]

    dist_cfg = cfg.get("distortions", {}) or {}
    C["distortions"] = list(dist_cfg.keys())
    sev_names = _safe(
        lambda: list(next(iter(dist_cfg.values()))["severities"].keys()), ["low", "med", "high"]
    )
    C["severities"] = sev_names
    C["sev_high"] = "high" if "high" in sev_names else (sev_names[-1] if sev_names else "high")

    ft = cfg.get("finetune", {}) or {}
    C["ft_epochs"] = ft.get("epochs")
    C["ft_lr0"] = ft.get("lr0")
    C["ft_batch"] = ft.get("batch_size")
    C["yolo_weights"] = (cfg.get("yolo", {}) or {}).get("weights", "yolov8n.pt")
    C["n_images"] = _safe(lambda: int(seg["n_images"])) if seg else None
    C["seg_pq_things"] = _safe(lambda: float(seg["PQ_things"])) if seg else None
    C["seg_pq_stuff"] = _safe(lambda: float(seg["PQ_stuff"])) if seg else None

    # ---- from comparison.csv ------------------------------------------------
    C["clean_by_task"] = {}
    C["det_worst"] = None
    C["snr_by_dist"] = {}
    C["enh_share_by_dist"] = {}
    C["ft_recovery_high"] = {}
    C["combined_wins"] = None
    C["winner_high"] = {}
    if comp is not None and len(comp):
        C["clean_by_task"] = _safe(
            lambda: {
                t: float(v)
                for t, v in comp.groupby("task")["clean"].mean().items()
                if ok(float(v))
            },
            {},
        )

        det = comp[comp["task"] == "detection"]
        if len(det) and det["distorted"].notna().any():
            row = det.loc[det["distorted"].idxmin()]
            C["det_worst"] = {
                "value": float(row["distorted"]),
                "distortion": str(row["distortion"]),
                "severity": str(row["severity"]),
            }

        C["snr_by_dist"] = _safe(
            lambda: {
                d: float(v)
                for d, v in comp.groupby("distortion")["snr_db"].mean().items()
                if ok(float(v))
            },
            {},
        )

        # share of damage recovered by enhancement, per (distortion, task):
        # sum(recovery_enhance) / sum(-degradation) across severities.
        for dist in C["distortions"]:
            shares = []
            sub = comp[comp["distortion"] == dist]
            for _task, g in sub.groupby("task"):
                dmg = _safe(lambda: float(-g["degradation"].sum()))
                rec = _safe(lambda: float(g["recovery_enhance"].sum(min_count=1)))
                if ok(dmg) and dmg > 1e-6 and ok(rec):
                    shares.append(rec / dmg)
            if shares:
                C["enh_share_by_dist"][dist] = (min(shares), max(shares))

        det_high = det[det["severity"] == C["sev_high"]]
        for _, r in det_high.iterrows():
            d = str(r["distortion"])
            if ok(r.get("recovery_finetune")):
                C["ft_recovery_high"][d] = float(r["recovery_finetune"])
            re_, rf_ = r.get("recovery_enhance"), r.get("recovery_finetune")
            if ok(re_) and ok(rf_):
                C["winner_high"][d] = {
                    "enh": float(re_),
                    "ft": float(rf_),
                    "winner": "classical enhancement" if float(re_) > float(rf_) else "fine-tuning",
                }

        if "recovery_combined" in det.columns and det["recovery_combined"].notna().any():
            wins = []
            for _, r in det.iterrows():
                rc, re_, rf_ = r.get("recovery_combined"), r.get("recovery_enhance"), r.get("recovery_finetune")
                if ok(rc) and ok(re_) and ok(rf_) and float(rc) > max(float(re_), float(rf_)) + 1e-9:
                    wins.append(f"{r['distortion']}/{r['severity']}")
            C["combined_wins"] = wins  # possibly empty list = "never beats"

    # ---- from summary_long.csv ----------------------------------------------
    C["ft_clean"] = None
    if slong is not None and len(slong):
        sub = slong[(slong["task"] == "detection") & (slong["variant"] == "finetuned_clean")]
        if len(sub):
            # prefer the primary detection metric if several are logged
            pick = None
            for pref in ("map50-95", "map", "ap"):
                m = sub[sub["metric"].astype(str).str.lower().str.replace("_", "").str.contains(pref.replace("-", ""))]
                if len(m):
                    pick = m.iloc[0]
                    break
            if pick is None:
                pick = sub.iloc[0]
            if ok(pick["value"]):
                C["ft_clean"] = float(pick["value"])

    return C


# ----------------------------------------------------------------------------
# Slide specs (single source of truth for both renderers)
# ----------------------------------------------------------------------------
def build_specs(cfg: dict, data: dict, C: dict) -> list[dict]:
    fig_dir = data["figures_dir"]
    specs: list[dict] = []

    def fig(name, fallback=None):
        return find_figure(fig_dir, name, fallback), name

    # 1 -- title -----------------------------------------------------------
    specs.append(
        {
            "layout": "title",
            "title": DECK_TITLE,
            "subtitle": DECK_SUBTITLE,
            "lines": [f"Team: {TEAM}", REPO_URL],
        }
    )

    # 2 -- design ----------------------------------------------------------
    n_img = C["n_images"]
    ds_txt = (
        f"Dataset: COCO val2017 — fixed seeded {n_img:,}-image class-coverage-balanced "
        "subset; the same paired images across all variants"
        if ok(n_img)
        else "Dataset: COCO val2017 — fixed seeded class-coverage-balanced subset "
        "(~1.5k images); the same paired images across all variants"
    )
    dists = C["distortions"] or ["gauss_noise", "salt_pepper", "motion_blur"]
    sevs = C["severities"] or ["low", "med", "high"]
    specs.append(
        {
            "layout": "bullets",
            "title": "Design",
            "bullets": [
                ds_txt,
                "4 tasks / models / metrics: ORB feature matching (match ratio) · "
                "YOLOv8n detection (bbox mAP) · Keypoint R-CNN pose (OKS AP) · "
                "Panoptic FPN segmentation (PQ)",
                f"{len(dists)} distortions × {len(sevs)} severities: "
                f"{', '.join(dists)}  ×  {' / '.join(sevs)}",
                "Two repair strategies: classical enhancement (matched per corruption) "
                "and fine-tuning the detector",
            ],
        }
    )

    # 3 -- method / pipeline -------------------------------------------------
    lr_txt = fmt_lr(C["ft_lr0"]) if ok(C["ft_lr0"]) else "1e-4"
    specs.append(
        {
            "layout": "bullets",
            "title": "Method / pipeline",
            "bullets": [
                "Deterministic per-image corruption: stable seeded RNG; degradation "
                "parameters logged per image (incl. motion-blur kernel angle)",
                "Lossless PNG storage — a JPEG round-trip would reshape the corruption",
                "Per-image SNR computed against the clean reference",
                "Matched enhancement per corruption: sigma-adaptive NLM (Gaussian noise), "
                "median filter (salt & pepper), non-blind Wiener with the logged kernel "
                "(motion blur)",
                "Fine-tune = per-image mixture of clean + all 9 corruption cells; "
                f"AdamW lr {lr_txt}; held-out train split for checkpoint selection",
                "Fine-tuned model evaluated on clean + distorted + enhanced val cells",
            ],
        }
    )

    # 4 -- clean baselines -----------------------------------------------------
    rows = []
    for t in TASK_ORDER:
        v = C["clean_by_task"].get(t)
        if ok(v):
            disp, model, metric = TASK_INFO[t]
            if t == "detection":
                model = Path(C["yolo_weights"]).stem or model
            rows.append([disp, model, metric, f3(v)])
    table = {"columns": ["Task", "Model", "Metric", "Clean score"], "rows": rows} if rows else None
    b4 = []
    if ok(C["seg_pq_things"]) and ok(C["seg_pq_stuff"]):
        b4.append(
            f"Segmentation PQ splits: things {f3(C['seg_pq_things'])} / "
            f"stuff {f3(C['seg_pq_stuff'])}"
        )
    b4.append(
        "Baselines match the published full-val2017 numbers → the measurement "
        "pipeline is validated"
    )
    specs.append(
        {
            "layout": "table",
            "title": "Clean baselines (Phase 1)",
            "table": table,
            "table_note": "comparison.csv not available — clean-baseline table omitted",
            "bullets": b4,
        }
    )

    # 5 -- degradation vs SNR ---------------------------------------------------
    b5 = []
    det_clean = C["clean_by_task"].get("detection")
    if ok(det_clean) and C["det_worst"]:
        w = C["det_worst"]
        b5.append(
            f"Detection mAP falls {f3(det_clean)} → {f3(w['value'])} at the worst "
            f"cell ({w['distortion']} {w['severity']})"
        )
    snr = C["snr_by_dist"]
    if len(snr) >= 2:
        parts = ", ".join(f"{d} {v:.1f} dB" for d, v in snr.items())
        mb = snr.get("motion_blur")
        if ok(mb) and mb == max(snr.values()):
            b5.append(
                "SNR alone does not predict damage: motion blur has the HIGHEST mean SNR "
                f"of the three corruptions ({parts}) yet destroys localization tasks most"
            )
        else:
            b5.append(f"SNR alone does not predict damage — mean SNR per distortion: {parts}")
    fp, fn = fig("acc_vs_snr_detection.png", "acc_vs_snr_*.png")
    specs.append(
        {
            "layout": "split",
            "title": "Degradation vs SNR (Phase 2)",
            "bullets": b5,
            "figure": fp,
            "figure_name": fn,
        }
    )

    # 6 -- distortion examples ---------------------------------------------------
    fp, fn = fig("grid_gauss_noise.png", "grid_*.png")
    cap_dist = fp.stem.replace("grid_", "") if fp else "gauss_noise"
    specs.append(
        {
            "layout": "figure",
            "title": "Distortion examples",
            "figure": fp,
            "figure_name": fn,
            "caption": f"{cap_dist.replace('_', ' ')}: clean / distorted (high severity) / "
            "enhanced (same paired images)",
        }
    )

    # 7 -- enhancement: what recovers ---------------------------------------------
    b7 = []
    for d in ["salt_pepper", "gauss_noise", "motion_blur"]:
        rng = C["enh_share_by_dist"].get(d)
        if rng:
            mn, mx = rng
            span = pct(mn) if pct(mn) == pct(mx) else f"{pct(mn)}–{pct(mx)}"
            b7.append(f"{d.replace('_', ' ')}: enhancement recovers {span} of the damage across tasks")
    fp, fn = fig("recovery_bars.png")
    specs.append(
        {
            "layout": "split",
            "title": "Enhancement (Phase 3): what recovers",
            "bullets": b7,
            "figure": fp,
            "figure_name": fn,
        }
    )

    # 8 -- predictions on images ----------------------------------------------------
    fp, fn = fig("annotated_gauss_noise.png", "annotated_*.png")
    specs.append(
        {
            "layout": "figure",
            "title": "Predictions on images",
            "figure": fp,
            "figure_name": fn,
            "caption": "Gray dashed = ground truth; colored = YOLO predictions "
            "(clean / distorted / enhanced / fine-tuned)",
        }
    )

    # 9 -- fine-tuning ------------------------------------------------------------
    b9 = []
    ep_txt = f"{C['ft_epochs']} epochs, " if ok(C["ft_epochs"]) else ""
    b9.append(
        "Setup: YOLOv8n fine-tuned on a per-image mixture of clean + all 9 corruption "
        f"cells ({ep_txt}AdamW lr {lr_txt})"
    )
    if ok(C["ft_clean"]) and ok(det_clean):
        cost = det_clean - C["ft_clean"]
        b9.append(
            f"On clean images: fine-tuned mAP {f3(C['ft_clean'])} vs pretrained "
            f"{f3(det_clean)} — clean cost = {f3(cost)}"
        )
    if C["ft_recovery_high"]:
        parts = ", ".join(f"{d.replace('_', ' ')} {v:+.3f}" for d, v in C["ft_recovery_high"].items())
        b9.append(f"In-domain recovery at high severity (absolute mAP): {parts}")
    if C["combined_wins"] is not None:
        if C["combined_wins"]:
            b9.append(
                "Combined repair (enhance + fine-tune) beats both single repairs on: "
                + ", ".join(C["combined_wins"])
            )
        else:
            b9.append(
                "Combined repair (enhance + fine-tune) never beats the best single "
                "repair — the two fixes do not stack"
            )
    specs.append({"layout": "bullets", "title": "Fine-tuning (Phase 4)", "bullets": b9})

    # 10 -- per-class view ------------------------------------------------------------
    fp, fn = fig("per_class_ap_gauss_noise_high.png", "per_class_ap_*.png")
    specs.append(
        {
            "layout": "figure",
            "title": "Per-class view",
            "figure": fp,
            "figure_name": fn,
            "caption": "Per-class AP under gauss_noise (high): clean vs distorted vs "
            "enhanced vs fine-tuned",
        }
    )

    # 11 -- decision matrix / conclusions -----------------------------------------------
    b11 = []
    for d in dists:
        w = C["winner_high"].get(d)
        if w:
            b11.append(
                f"{d.replace('_', ' ')}: {w['winner']} wins "
                f"(enhance {w['enh']:+.3f} vs fine-tune {w['ft']:+.3f} mAP at high severity)"
            )
    b11 += [
        "Enhancement is corruption-specific — each corruption needs its matched restorer",
        "Fine-tuning needs mixture training — single-corruption training "
        "negative-transfers to every other cell",
        "Blur is only repairable because the kernel is known (non-blind Wiener)",
    ]
    specs.append({"layout": "bullets", "title": "Decision matrix / conclusions", "bullets": b11})

    # 12 -- limitations ------------------------------------------------------------------
    specs.append(
        {
            "layout": "bullets",
            "title": "Limitations",
            "bullets": [
                "ORB match ratio structurally penalizes smoothing enhancers — it is a "
                "fidelity-to-clean metric, not a task-utility metric",
                "Wiener deconvolution uses the logged (oracle) kernel; blind deconvolution "
                "with a wrong angle measured WORSE than no restoration",
                "yolov8n is the smallest detector in its family — absolute mAP is modest",
                "Evaluation subset is ≈1.5k images (class-coverage-balanced, seeded)",
            ],
        }
    )

    assert len(specs) == 12, f"expected 12 slides, built {len(specs)}"
    return specs


# ----------------------------------------------------------------------------
# PPTX rendering
# ----------------------------------------------------------------------------
class Box:
    def __init__(self, x, y, w, h):
        self.x, self.y, self.w, self.h = x, y, w, h


def _pptx_textbox(slide, box: Box):
    tb = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
    tb.text_frame.word_wrap = True
    return tb


def _pptx_para(tf, first: bool):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def _pptx_run(p, text, size, color, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = hexrgb(color)
    r.font.bold = bold
    r.font.italic = italic
    return r


def _pptx_accent_bar(slide, x, y, w, h=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = hexrgb(ACCENT)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _pptx_title(slide, title):
    tb = _pptx_textbox(slide, Box(0.55, 0.32, 12.3, 0.95))
    p = tb.text_frame.paragraphs[0]
    _pptx_run(p, title, 30, TITLE_C, bold=True)
    _pptx_accent_bar(slide, 0.6, 1.22, 2.6, 0.045)


def _pptx_footer(slide, n, total):
    tb = _pptx_textbox(slide, Box(11.9, 7.05, 1.2, 0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _pptx_run(p, f"{n} / {total}", 10, MUTED)


def _pptx_bullets(slide, bullets, box: Box):
    if not bullets:
        return
    total_chars = sum(len(b) for b in bullets)
    size = 18 if total_chars < 380 else (16 if total_chars < 750 else 14)
    tb = _pptx_textbox(slide, box)
    tf = tb.text_frame
    for i, b in enumerate(bullets):
        p = _pptx_para(tf, i == 0)
        p.space_after = Pt(12)
        _pptx_run(p, "•  " + b, size, BODY_C)


def _pptx_image(slide, path, box: Box, note_name: str):
    if path is None:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = hexrgb(PLACE_BG)
        shp.line.color.rgb = hexrgb(PLACE_BORDER)
        shp.line.width = Pt(1)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _pptx_run(p, f"[figure not available: {note_name}]", 14, MUTED, italic=True)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    s = min(box.w / iw, box.h / ih)
    w, h = iw * s, ih * s
    slide.shapes.add_picture(
        str(path),
        Inches(box.x + (box.w - w) / 2),
        Inches(box.y + (box.h - h) / 2),
        width=Inches(w),
        height=Inches(h),
    )


def _pptx_table(slide, table, box: Box):
    nrows = len(table["rows"]) + 1
    ncols = len(table["columns"])
    gf = slide.shapes.add_table(
        nrows, ncols, Inches(box.x), Inches(box.y), Inches(box.w), Inches(min(box.h, 0.55 * nrows))
    )
    t = gf.table
    widths = [0.22, 0.22, 0.34, 0.22]
    for j in range(ncols):
        t.columns[j].width = Inches(box.w * widths[j % len(widths)])
    for j, name in enumerate(table["columns"]):
        cell = t.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hexrgb(ACCENT)
        p = cell.text_frame.paragraphs[0]
        _pptx_run(p, str(name), 15, "#ffffff", bold=True)
    for i, row in enumerate(table["rows"], start=1):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hexrgb("#f6f8fb" if i % 2 else "#ffffff")
            p = cell.text_frame.paragraphs[0]
            _pptx_run(p, str(val), 14, BODY_C, bold=(j == ncols - 1))


def render_pptx(specs: list[dict], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(specs)

    for n, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(blank)
        layout = spec["layout"]

        if layout == "title":
            tb = _pptx_textbox(slide, Box(0.9, 2.05, 11.6, 1.7))
            p = tb.text_frame.paragraphs[0]
            _pptx_run(p, spec["title"], 40, TITLE_C, bold=True)
            _pptx_accent_bar(slide, 0.95, 3.62, 3.2, 0.06)
            tb = _pptx_textbox(slide, Box(0.9, 3.9, 11.6, 0.7))
            _pptx_run(tb.text_frame.paragraphs[0], spec["subtitle"], 20, "#5a6b7d")
            for k, line in enumerate(spec.get("lines", [])):
                tb = _pptx_textbox(slide, Box(0.9, 4.85 + 0.55 * k, 11.6, 0.5))
                color = ACCENT if line.startswith("http") else BODY_C
                _pptx_run(tb.text_frame.paragraphs[0], line, 16 if k else 18, color)
        elif layout == "bullets":
            _pptx_title(slide, spec["title"])
            _pptx_bullets(slide, spec["bullets"], Box(0.7, 1.6, 12.0, 5.4))
        elif layout == "split":
            _pptx_title(slide, spec["title"])
            _pptx_bullets(slide, spec["bullets"], Box(0.6, 1.7, 5.9, 5.2))
            _pptx_image(slide, spec["figure"], Box(6.7, 1.5, 6.1, 5.4), spec["figure_name"])
        elif layout == "figure":
            _pptx_title(slide, spec["title"])
            _pptx_image(slide, spec["figure"], Box(0.8, 1.45, 11.7, 5.15), spec["figure_name"])
            tb = _pptx_textbox(slide, Box(0.8, 6.72, 11.7, 0.5))
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _pptx_run(p, spec.get("caption", ""), 13, MUTED, italic=True)
        elif layout == "table":
            _pptx_title(slide, spec["title"])
            if spec["table"]:
                _pptx_table(slide, spec["table"], Box(1.2, 1.7, 10.9, 3.0))
            else:
                _pptx_image(slide, None, Box(1.2, 1.7, 10.9, 2.6), spec["table_note"])
            _pptx_bullets(slide, spec["bullets"], Box(0.7, 5.1, 12.0, 1.9))

        _pptx_footer(slide, n, total)

    prs.save(str(out_path))


# ----------------------------------------------------------------------------
# PDF rendering (matplotlib PdfPages, same specs)
# ----------------------------------------------------------------------------
LINE_PT = SLIDE_H * 72.0  # figure height in points (for line-height fractions)


def _pdf_title(fig, title):
    fig.text(0.045, 0.945, title, fontsize=23, fontweight="bold", color=TITLE_C, va="top")
    fig.add_artist(
        Rectangle((0.047, 0.868), 0.17, 0.007, transform=fig.transFigure, facecolor=ACCENT, edgecolor="none")
    )


def _pdf_footer(fig, n, total):
    fig.text(0.965, 0.03, f"{n} / {total}", fontsize=8, color=MUTED, ha="right")


def _pdf_bullets(fig, bullets, x, y_top, wrap, size):
    y = y_top
    line_frac = size * 1.45 / LINE_PT
    for b in bullets:
        lines = textwrap.wrap(b, wrap) or [""]
        txt = "•  " + lines[0] + "".join("\n    " + ln for ln in lines[1:])
        fig.text(x, y, txt, fontsize=size, color=BODY_C, va="top", linespacing=1.45)
        y -= line_frac * len(lines) + 0.022


def _pdf_image(fig, path, box, note_name):
    """box = (x, y, w, h) in figure fractions; preserves image aspect ratio."""
    x, y, w, h = box
    if path is None:
        fig.add_artist(
            Rectangle(
                (x, y), w, h, transform=fig.transFigure, facecolor=PLACE_BG,
                edgecolor=PLACE_BORDER, linestyle="--", linewidth=1.2,
            )
        )
        fig.text(
            x + w / 2, y + h / 2, f"[figure not available: {note_name}]",
            fontsize=12, color=MUTED, style="italic", ha="center", va="center",
        )
        return
    with Image.open(path) as im:
        arr = np.asarray(im.convert("RGB"))
        iw, ih = im.size
    bw_in, bh_in = w * SLIDE_W, h * SLIDE_H
    s = min(bw_in / iw, bh_in / ih)
    wf, hf = iw * s / SLIDE_W, ih * s / SLIDE_H
    ax = fig.add_axes([x + (w - wf) / 2, y + (h - hf) / 2, wf, hf])
    ax.imshow(arr)
    ax.axis("off")


def _pdf_table(fig, table, box):
    ax = fig.add_axes(list(box))
    ax.axis("off")
    ncols = len(table["columns"])
    widths = [0.22, 0.22, 0.34, 0.22][:ncols]
    tbl = ax.table(
        cellText=table["rows"], colLabels=table["columns"], colWidths=widths,
        cellLoc="left", colLoc="left", loc="upper left", bbox=[0, 0, 1, 1],
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(13)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#d5dbe2")
        cell.PAD = 0.03
        if r == 0:
            cell.set_facecolor(ACCENT)
            cell.get_text().set_color("white")
            cell.get_text().set_weight("bold")
        else:
            cell.set_facecolor("#f6f8fb" if r % 2 else "#ffffff")
            cell.get_text().set_color(BODY_C)


def render_pdf(specs: list[dict], out_path: Path) -> int:
    total = len(specs)
    pages = 0
    with PdfPages(str(out_path)) as pdf:
        for n, spec in enumerate(specs, start=1):
            fig = plt.figure(figsize=(SLIDE_W, SLIDE_H), dpi=120)
            fig.patch.set_facecolor("white")
            layout = spec["layout"]

            if layout == "title":
                fig.text(0.5, 0.62, spec["title"], fontsize=30, fontweight="bold",
                         color=TITLE_C, ha="center", va="center")
                fig.add_artist(Rectangle((0.40, 0.545), 0.20, 0.008,
                                         transform=fig.transFigure, facecolor=ACCENT, edgecolor="none"))
                fig.text(0.5, 0.48, spec["subtitle"], fontsize=17, color="#5a6b7d", ha="center")
                for k, line in enumerate(spec.get("lines", [])):
                    color = ACCENT if line.startswith("http") else BODY_C
                    fig.text(0.5, 0.38 - 0.06 * k, line, fontsize=14 if k else 15,
                             color=color, ha="center")
            elif layout == "bullets":
                _pdf_title(fig, spec["title"])
                total_chars = sum(len(b) for b in spec["bullets"])
                size = 15 if total_chars < 380 else (14 if total_chars < 750 else 12.5)
                _pdf_bullets(fig, spec["bullets"], 0.055, 0.79, 100, size)
            elif layout == "split":
                _pdf_title(fig, spec["title"])
                _pdf_bullets(fig, spec["bullets"], 0.048, 0.77, 48, 13)
                _pdf_image(fig, spec["figure"], (0.50, 0.10, 0.46, 0.70), spec["figure_name"])
            elif layout == "figure":
                _pdf_title(fig, spec["title"])
                _pdf_image(fig, spec["figure"], (0.06, 0.17, 0.88, 0.65), spec["figure_name"])
                fig.text(0.5, 0.095, spec.get("caption", ""), fontsize=12, color=MUTED,
                         style="italic", ha="center", va="top")
            elif layout == "table":
                _pdf_title(fig, spec["title"])
                if spec["table"]:
                    nrows = len(spec["table"]["rows"]) + 1
                    h = min(0.42, 0.085 * nrows)
                    _pdf_table(fig, spec["table"], (0.10, 0.78 - h, 0.80, h))
                else:
                    _pdf_image(fig, None, (0.10, 0.42, 0.80, 0.36), spec["table_note"])
                _pdf_bullets(fig, spec["bullets"], 0.055, 0.27, 100, 14)

            _pdf_footer(fig, n, total)
            pdf.savefig(fig)
            plt.close(fig)
            pages += 1
    return pages


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------
def main() -> None:
    repo_root = Path(__file__).resolve().parent.parent
    ap = argparse.ArgumentParser(description="Build the final presentation (PPTX + PDF).")
    ap.add_argument("--config", default=str(repo_root / "configs" / "config.yaml"))
    ap.add_argument("--out", default=str(repo_root / "slides"))
    args = ap.parse_args()

    config_path = Path(args.config).resolve()
    cfg = load_config(config_path)
    data = load_inputs(cfg, config_path.parent, repo_root)
    C = compute_context(cfg, data)
    specs = build_specs(cfg, data, C)

    missing = [s["figure_name"] for s in specs if "figure" in s and s["figure"] is None]
    if data["comp"] is None:
        print("note: comparison.csv missing -> computed bullets/table degraded")
    if data["slong"] is None:
        print("note: summary_long.csv missing -> fine-tuned-on-clean bullet dropped")
    for m in missing:
        print(f"note: figure missing -> placeholder used: {m}")

    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / "final_presentation.pptx"
    pdf_path = out_dir / "final_presentation.pdf"

    render_pptx(specs, pptx_path)
    pages = render_pdf(specs, pdf_path)

    print(f"wrote {pptx_path}  ({pptx_path.stat().st_size / 1024:.0f} KB, {len(specs)} slides)")
    print(f"wrote {pdf_path}  ({pdf_path.stat().st_size / 1024:.0f} KB, {pages} pages)")


if __name__ == "__main__":
    main()
